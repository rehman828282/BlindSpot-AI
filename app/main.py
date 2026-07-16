import io
import json
import math
import hashlib
import sqlite3
from pathlib import Path
from typing import Any

from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel

from app.debate_logic import draft_missing_piece, stream_debate

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None


BASE_DIR = Path(__file__).resolve().parent.parent
TEMPLATE_DIR = BASE_DIR / "templates"
SQLITE_PATH = BASE_DIR / "blindspot_store.sqlite3"
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx"}

app = FastAPI(title="BlindSpot AI", version="0.1.0")
templates = Jinja2Templates(directory=str(TEMPLATE_DIR))


class FixRequest(BaseModel):
    issue: str
    instructions: str | None = None


def cosine_similarity(left: list[float], right: list[float]) -> float:
    return sum(a * b for a, b in zip(left, right))


class SQLiteVectorStore:
    def __init__(self) -> None:
        self._init_sqlite()

    def _connect(self) -> sqlite3.Connection:
        return sqlite3.connect(SQLITE_PATH)

    def _init_sqlite(self) -> None:
        with self._connect() as connection:
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS document_chunks (
                    id TEXT PRIMARY KEY,
                    document TEXT NOT NULL,
                    source TEXT NOT NULL,
                    chunk_index INTEGER NOT NULL,
                    embedding TEXT NOT NULL
                )
                """
            )

    def reset(self) -> None:
        self._init_sqlite()
        with self._connect() as connection:
            connection.execute("DELETE FROM document_chunks")

    def add(self, *, chunks: list[str], filename: str) -> None:
        ids = [f"chunk-{index}" for index in range(len(chunks))]
        embeddings = [embed_text(chunk) for chunk in chunks]

        with self._connect() as connection:
            connection.executemany(
                """
                INSERT INTO document_chunks (id, document, source, chunk_index, embedding)
                VALUES (?, ?, ?, ?, ?)
                """,
                [
                    (
                        ids[index],
                        chunk,
                        filename,
                        index,
                        json.dumps(embeddings[index], ensure_ascii=True),
                    )
                    for index, chunk in enumerate(chunks)
                ],
            )

    def count(self) -> int:
        self._init_sqlite()
        with self._connect() as connection:
            row = connection.execute("SELECT COUNT(*) FROM document_chunks").fetchone()
        return int(row[0]) if row else 0

    def query(self, query_embedding: list[float], *, limit: int = 12) -> list[tuple[str, dict[str, Any]]]:
        count = self.count()
        if count == 0:
            return []

        self._init_sqlite()
        with self._connect() as connection:
            rows = connection.execute(
                "SELECT document, source, chunk_index, embedding FROM document_chunks"
            ).fetchall()

        ranked = []
        for document, source, chunk_index, embedding_json in rows:
            embedding = json.loads(embedding_json)
            ranked.append(
                (
                    cosine_similarity(query_embedding, embedding),
                    document,
                    {"source": source, "chunk": chunk_index},
                )
            )

        ranked.sort(key=lambda item: item[0], reverse=True)
        return [(document, metadata) for _, document, metadata in ranked[:limit]]


vector_store = SQLiteVectorStore()


def extract_pdf_text(data: bytes) -> str:
    if PdfReader is None:
        raise HTTPException(status_code=500, detail="PDF parsing requires pypdf. Install requirements.txt.")
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx_text(data: bytes) -> str:
    if Document is None:
        raise HTTPException(status_code=500, detail="DOCX parsing requires python-docx. Install requirements.txt.")
    document = Document(io.BytesIO(data))
    parts: list[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_pptx_text(data: bytes) -> str:
    if Presentation is None:
        raise HTTPException(status_code=500, detail="PPTX parsing requires python-pptx. Install requirements.txt.")
    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        slides.append("\n".join(lines))

    return "\n\n".join(slides)


def extract_text(data: bytes, extension: str) -> str:
    if extension == ".pdf":
        return extract_pdf_text(data)
    if extension == ".docx":
        return extract_docx_text(data)
    if extension == ".pptx":
        return extract_pptx_text(data)
    raise HTTPException(status_code=400, detail="Unsupported file type.")


def chunk_text(text: str, *, chunk_words: int = 700, overlap_words: int = 100) -> list[str]:
    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    step = max(1, chunk_words - overlap_words)
    for start in range(0, len(words), step):
        chunk = " ".join(words[start : start + chunk_words])
        if chunk.strip():
            chunks.append(chunk)
        if start + chunk_words >= len(words):
            break
    return chunks


def embed_text(text: str, *, dimensions: int = 384) -> list[float]:
    vector = [0.0] * dimensions
    tokens = text.lower().split()

    if not tokens:
        vector[0] = 1.0
        return vector

    for token in tokens:
        digest = hashlib.blake2b(token.encode("utf-8"), digest_size=8).digest()
        index = int.from_bytes(digest[:4], "little") % dimensions
        sign = 1.0 if digest[4] % 2 == 0 else -1.0
        vector[index] += sign

    magnitude = math.sqrt(sum(value * value for value in vector)) or 1.0
    return [value / magnitude for value in vector]


def query_context() -> str:
    count = vector_store.count()
    if count == 0:
        return ""

    query = (
        "critical missing piece unresolved decision unsupported assumptions evidence risk "
        "user needs competitor standards compliance"
    )
    context_blocks = []
    for doc, metadata in vector_store.query(embed_text(query), limit=min(count, 12)):
        label = metadata.get("source", "uploaded document") if isinstance(metadata, dict) else "uploaded document"
        context_blocks.append(f"[{label}]\n{doc}")

    return "\n\n---\n\n".join(context_blocks)


def sse(event: str, data: Any) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=True)}\n\n"


@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    filename = file.filename or "uploaded-document"
    extension = Path(filename).suffix.lower()

    if extension not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Upload a .pdf, .docx, or .pptx file.")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="The uploaded file was empty.")

    try:
        text = extract_text(data, extension)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not parse document: {exc}") from exc

    cleaned_text = text.strip()
    if not cleaned_text:
        raise HTTPException(status_code=400, detail="No extractable text was found in the uploaded file.")

    chunks = chunk_text(cleaned_text)
    vector_store.reset()
    vector_store.add(chunks=chunks, filename=filename)

    return JSONResponse(
        {
            "filename": filename,
            "characters": len(cleaned_text),
            "chunks": len(chunks),
            "store": "sqlite",
            "preview": cleaned_text[:6000],
        }
    )


@app.get("/analyze")
async def analyze():
    async def event_stream():
        context = query_context()
        if not context:
            yield sse("error", {"message": "Upload a document before starting analysis."})
            yield sse("done", {"message": "No analysis was run."})
            return

        try:
            async for update in stream_debate(context):
                yield sse(update["event"], update["data"])
        except Exception as exc:
            yield sse("error", {"message": str(exc)})

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/fix")
async def fix(payload: FixRequest):
    issue = payload.issue.strip()
    if not issue:
        raise HTTPException(status_code=400, detail="A critical issue is required.")

    result = await draft_missing_piece(
        issue=issue,
        instructions=payload.instructions,
        context=query_context(),
    )
    return JSONResponse(result)
